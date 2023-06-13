import os
import subprocess
import configparser
from pptx import Presentation

def insert_names_into_ppt(names_file, template_ppt, output_folder, libreoffice_path):
    # Load the names from the text file
    with open(names_file, 'r') as file:
        names = file.readlines()

    # Iterate over each name and create a modified copy of the presentation
    for name in names:
        # Load the PowerPoint template for each name
        prs = Presentation(template_ppt)

        # Iterate over each slide in the presentation and replace placeholder text
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if '<<name>>' in run.text:
                                run.text = run.text.replace('<<name>>', name.strip())

        # Save the modified PowerPoint presentation with the name
        output_pptx = f"{output_folder}/{name.strip()}.pptx"
        prs.save(output_pptx)

        # Convert the PowerPoint file to PDF using LibreOffice
        output_pdf = f"{output_folder}/{name.strip()}.pdf"
        subprocess.call([libreoffice_path, '--headless', '--convert-to', 'pdf', '--outdir', output_folder, output_pptx])

        # Remove the intermediate PowerPoint file
        os.remove(output_pptx)

# Load the configuration from the config file
config = configparser.ConfigParser()
config.read('config.ini')

# Get the values from the config file
libreoffice_path = config.get('Paths', 'LibreOfficePath')
names_file = config.get('Paths', 'NamesFile')
template_pptx = config.get('Paths', 'TemplatePPTX')
output_folder = config.get('Paths', 'OutputFolder')

# Usage example
insert_names_into_ppt(names_file, template_pptx, output_folder, libreoffice_path)
